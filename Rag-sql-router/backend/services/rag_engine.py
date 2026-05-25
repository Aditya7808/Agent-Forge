import os
import uuid
from pathlib import Path
from openai import OpenAI
import chromadb
from chromadb.utils import embedding_functions
from backend.config import settings

chroma_client = chromadb.PersistentClient(path=settings.chroma_persist_dir)
sessions: dict = {}


def get_openai_ef():
    return embedding_functions.OpenAIEmbeddingFunction(
        api_key=settings.openai_api_key,
        model_name=settings.embedding_model,
    )


def extract_text_from_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    elif ext == ".pdf":
        try:
            import PyPDF2
            text_parts = []
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text_parts.append(page.extract_text() or "")
            return "\n".join(text_parts)
        except ImportError:
            raise ImportError("PyPDF2 is required for PDF processing. Install with: pip install PyPDF2")
    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except ImportError:
            raise ImportError("python-docx is required for DOCX processing. Install with: pip install python-docx")
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except ImportError:
            raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
    else:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 200) -> list:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks


def process_documents(session_id: str, file_dir: str) -> int:
    collection_name = f"session_{session_id.replace('-', '_')[:32]}"

    openai_ef = get_openai_ef()
    collection = chroma_client.get_or_create_collection(
        name=collection_name,
        embedding_function=openai_ef,
        metadata={"hnsw:space": "cosine"}
    )

    total_chunks = 0
    for file_name in os.listdir(file_dir):
        file_path = os.path.join(file_dir, file_name)
        if not os.path.isfile(file_path):
            continue

        text = extract_text_from_file(file_path)
        if not text.strip():
            continue

        chunks = chunk_text(text)
        ids = [f"{file_name}_{i}_{uuid.uuid4().hex[:8]}" for i in range(len(chunks))]
        metadatas = [{"source": file_name, "chunk_index": i} for i in range(len(chunks))]

        batch_size = 100
        for i in range(0, len(chunks), batch_size):
            collection.add(
                documents=chunks[i:i+batch_size],
                ids=ids[i:i+batch_size],
                metadatas=metadatas[i:i+batch_size],
            )

        total_chunks += len(chunks)

    sessions[session_id] = collection_name
    return total_chunks


def query_documents(session_id: str, query: str) -> dict:
    collection_name = sessions.get(session_id)
    if not collection_name:
        return {
            "response": "No documents have been uploaded for this session. Please upload documents first.",
            "sources": [],
            "context": ""
        }

    openai_ef = get_openai_ef()
    collection = chroma_client.get_collection(
        name=collection_name,
        embedding_function=openai_ef,
    )

    results = collection.query(
        query_texts=[query],
        n_results=settings.similarity_top_k,
    )

    if not results["documents"][0]:
        return {
            "response": "No relevant information found in the uploaded documents.",
            "sources": [],
            "context": ""
        }

    context_parts = results["documents"][0]
    sources = [m["source"] for m in results["metadatas"][0]]
    context_str = "\n\n---\n\n".join(context_parts)

    client = OpenAI(api_key=settings.openai_api_key)

    response = client.chat.completions.create(
        model=settings.openai_model,
        messages=[
            {"role": "system", "content": (
                "You are a meticulous document analyst. Answer the user's question based EXCLUSIVELY on the provided context. "
                "Rules:\n"
                "1. Ground your entire response in the facts from the context. Do not use prior knowledge.\n"
                "2. If multiple parts are relevant, synthesize them into a coherent answer.\n"
                "3. If the context lacks sufficient information, state: 'The provided documents do not contain enough information to answer this question.'\n"
                "4. Cite sources when possible."
            )},
            {"role": "user", "content": f"Context:\n{context_str}\n\n---\n\nQuestion: {query}"}
        ],
        temperature=0.2,
        max_tokens=1000,
    )

    return {
        "response": response.choices[0].message.content.strip(),
        "sources": list(set(sources)),
        "context": context_str
    }
