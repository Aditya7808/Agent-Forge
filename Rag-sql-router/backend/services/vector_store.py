import os
import logging
import shutil
from pathlib import Path
from typing import List
from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from backend.config import settings
from backend.services.llm import get_embeddings

logger = logging.getLogger(__name__)

_sessions: dict[str, str] = {}


def _session_collection_name(session_id: str) -> str:
    safe = session_id.replace("-", "_")[:32]
    return f"session_{safe}"


def get_vector_store(session_id: str) -> Chroma:
    collection_name = _session_collection_name(session_id)
    return Chroma(
        collection_name=collection_name,
        embedding_function=get_embeddings(),
        persist_directory=settings.chroma_persist_dir,
    )


def has_documents(session_id: str) -> bool:
    return session_id in _sessions


def load_documents_from_dir(file_dir: str) -> List[Document]:
    docs: List[Document] = []
    for filename in os.listdir(file_dir):
        file_path = os.path.join(file_dir, filename)
        if not os.path.isfile(file_path):
            continue

        ext = Path(filename).suffix.lower()

        try:
            if ext == ".pdf":
                from langchain_community.document_loaders import PyPDFLoader
                loader = PyPDFLoader(file_path)
                file_docs = loader.load()
            elif ext == ".docx":
                from langchain_community.document_loaders import Docx2txtLoader
                loader = Docx2txtLoader(file_path)
                file_docs = loader.load()
            elif ext == ".pptx":
                from langchain_community.document_loaders import UnstructuredPowerPointLoader
                try:
                    loader = UnstructuredPowerPointLoader(file_path)
                    file_docs = loader.load()
                except Exception:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_parts.append(shape.text)
                    file_docs = [Document(
                        page_content="\n".join(text_parts),
                        metadata={"source": filename}
                    )]
            elif ext == ".txt":
                from langchain_community.document_loaders import TextLoader
                loader = TextLoader(file_path, encoding="utf-8", autodetect_encoding=True)
                file_docs = loader.load()
            else:
                logger.warning(f"Skipping unsupported file: {filename}")
                continue

            for d in file_docs:
                d.metadata["source"] = filename
            docs.extend(file_docs)
        except Exception as e:
            logger.exception(f"Failed to load {filename}: {e}")
            continue

    return docs


def ingest_documents(session_id: str, file_dir: str) -> int:
    docs = load_documents_from_dir(file_dir)
    if not docs:
        return 0

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=200,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    chunks = splitter.split_documents(docs)

    if not chunks:
        return 0

    vs = get_vector_store(session_id)
    vs.add_documents(chunks)
    _sessions[session_id] = _session_collection_name(session_id)

    logger.info(f"Ingested {len(chunks)} chunks for session {session_id}")
    return len(chunks)


def clear_session(session_id: str) -> None:
    try:
        vs = get_vector_store(session_id)
        vs.delete_collection()
    except Exception as e:
        logger.warning(f"Failed to delete collection: {e}")
    _sessions.pop(session_id, None)
